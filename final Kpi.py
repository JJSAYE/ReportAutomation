import os
os.chdir("/Users/jsaye/projects/technomic/GitHub/")
import sys

import datetime
import json
from shutil import copyfile
import xlrd
import math 

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.util import Pt
from PIL import Image

def char_range(c1, c2):
	for c in range(ord(c1), ord(c2)+1):
		yield chr(c)

def slugify(s):
	result = ""
	for c in s.lower():
		if c == " ":
			continue
		else:
			if c in char_range("a","z"):
				result += c 
			else:
				result += "_"
	return result 

def is_mac_not_windows():
	return sys.platform != "win32"

if is_mac_not_windows():
	AUTOMATION_PATH = "/Users/jsaye/projects/technomic/GitHub/"
	LOGO_PATH = AUTOMATION_PATH + "/logos/"
	DESTINATION_PATH = AUTOMATION_PATH + "finished_reports/"
else:
	AUTOMATION_PATH = "C:\\Users\\jsaye\\Documents\\Work\\Python Scripts\\Automation" + "\\" 
	LOGO_PATH = AUTOMATION_PATH + "logos" + "\\"
	DESTINATION_PATH = AUTOMATION_PATH + "\\finished reports\\"  

#Comp Set function
def comp_set(VA_stats, col, count_points, Chain):
    result = []
    all_chains = list(VA_stats.keys())
    all_chains.sort()
    for row_index in range(0, count_points):
            row_name = all_chains[row_index]
            data_for_chain = VA_stats[Chain]
            result.append(data_for_chain)
    return data_for_chain

#Extraction function
def extract_vector(all_stats, col, count_points):
    result = []
    all_chains = list(all_stats.keys())
    all_chains.sort()

    for row_index in range(0, count_points):
        row_name = all_chains[row_index]
        data_for_chain = all_stats[row_name]
        result.append(data_for_chain[col])
    return result

#Create chart function
def crave_chart(
				chart_name,
				slide_name,
				placeholder_index,
				categories,
				series
				):
	chart_data = ChartData()
	my_tuple = list(zip(categories, series))
	my_tuple.sort(key=lambda elem: elem[1])
	new_cat = [i[0] for i in my_tuple]
	new_series = [i[1] for i in my_tuple]
	chart_data.categories = (new_cat)
	chart_data.add_series('Series1', (new_series))
	chart_name = slide_name.placeholders[placeholder_index].insert_chart(XL_CHART_TYPE.BAR_CLUSTERED, chart_data)
	chart_name.chart.plots[0].series[0].format.fill.solid()
	chart_name.chart.plots[0].series[0].format.fill.fore_color.rgb = RGBColor(0,132,192)
	chart_name.chart.plots[0].has_data_labels = True
	chart_name.chart.value_axis.has_major_gridlines = False
	chart_name.chart.value_axis.visible = False
	chart_name.chart.plots[0].data_labels.font.size = Pt(14)
	chart_name.chart.plots[0].data_labels.number_format = '#,##0.0%'
	chart_name.chart.plots[0].chart.category_axis.format.line.fill.solid()
	chart_name.chart.plots[0].chart.category_axis.format.line.fill.fore_color.rgb = RGBColor(255,255,255)


def create_chart(
				chart_name, 
				slide_name, 
				placeholder_index, 
				categories, 
				series,
				name_of_attribute
				):
	chart_data = ChartData()
	my_tuple = list(zip(categories, series))
	my_tuple.sort(key=lambda elem: elem[1])
	new_cat = [i[0] for i in my_tuple]
	new_series = [i[1] for i in my_tuple]
	if  name_of_attribute == "Food taste and flavor":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Food taste and flavor"])
	elif name_of_attribute == "Overall Rating":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Overall Rating"])
	elif name_of_attribute == "Food quality":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Food quality"])
	elif name_of_attribute == "Food quality takeout":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Food quality takeout"])
	elif name_of_attribute == "Interior cleanliness":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Interior cleanliness"])
	elif name_of_attribute == "Kitchen/food prep area cleanliness":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Kitchen/food prep area cleanliness"])
	elif name_of_attribute == "Order accuracy":
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Order accuracy"])
	else:
		index = new_series.index(kpis[slugify(Chain)]["attributes"]["Dishware/glassware/silverware cleanliness"])

	# index = [name_of_attribute for name_of_attribute in new_series.index(kpis[slugify(Chain)]["attributes"][name_of_attribute]) if name_of_attribute == name_of_attribute]

	chart_data.categories = (new_cat)
	chart_data.add_series('Series1', (new_series))
	chart_name = slide_name.placeholders[placeholder_index].insert_chart(XL_CHART_TYPE.BAR_CLUSTERED, chart_data)
	chart_name.chart.plots[0].series[0].format.fill.solid()
	chart_name.chart.plots[0].series[0].format.fill.fore_color.rgb = RGBColor(192,192,192)
	chart_name.chart.plots[0].has_data_labels = True
	chart_name.chart.value_axis.has_major_gridlines = False
	chart_name.chart.value_axis.visible = False
	chart_name.chart.plots[0].data_labels.font.size = Pt(14)
	chart_name.chart.plots[0].data_labels.number_format = '#,##0.0%'
	chart_name.chart.plots[0].chart.value_axis.maximum_scale = .95
	chart_name.chart_part.chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
	point = chart_name.chart.plots[0].series[0].points[index]
	fill = point.format.fill
	fill.solid()
	fill.fore_color.rgb = RGBColor(225,34,34)
	chart_name.chart.plots[0].chart.category_axis.format.line.fill.solid()
	chart_name.chart.plots[0].chart.category_axis.format.line.fill.fore_color.rgb = RGBColor(255,255,255)

def logo_pic(placeholer_number, competitor_index_number):
	pic_name = vo_slide.placeholders[placeholer_number].insert_picture(LOGO_PATH 
																	+ kpis[va_slide_competitor_list[competitor_index_number]]["chain_name"]
																	+ '.png')
	pic_name.crop_right = 0
	pic_name.crop_left = 0
	pic_name.crop_top = 0
	pic_name.crop_bottom = 0

def write_text(slide_name, placeholder_number, text):
	slide_name.placeholders[placeholder_number].text = text

def write_base_text(slide_name, placeholder_number, text):
	slide_name.placeholders[placeholder_number].text = text
	slide_name.placeholders[placeholder_number].text_frame.paragraphs[0].font.size = Pt(8)
	slide_name.placeholders[placeholder_number].text_frame.paragraphs[0].font.fill.solid()
	slide_name.placeholders[placeholder_number].text_frame.paragraphs[0].font.fill.fore_color.rgb = RGBColor(137,141,141)
	slide_name.placeholders[placeholder_number].text_frame.paragraphs[0].font.name = "Arial (Body)"

def write_rest_visit_text(slide_name, placeholder_number, text, more_text):
	text_frame = slide_name.placeholders[placeholder_number].text_frame
	text_frame.paragraphs[0].font.size = Pt(20)
	text_frame.paragraphs[0].font.fill.solid()
	text_frame.paragraphs[0].font.fill.fore_color.rgb = RGBColor(137,141,141)
	text_frame.paragraphs[0].font.name = "Arial (Body)"
	p = text_frame.paragraphs[0]
	run = p.add_run()
	run.text = text
	run.font.bold = True
	run2 = p.add_run()
	run2.text = more_text

def dem_text(
			dem_type, 
			placeholder_index, 
			skew_number, 
			text_after_skew_number, 
			avg_number, 
			text_after_avg_number, 
			seg):
		text_frame = dem_slide.placeholders[placeholder_index].text_frame
		dem_type = dem_slide.placeholders[placeholder_index].text_frame
		skew = dem_type.paragraphs[0]
		run = skew.add_run()
		run.text = skew_number
		font = run.font
		font.bold = True
		font.color.rgb = RGBColor(200,34,34)
		text_after_skew = text_frame.paragraphs[0]
		run2 = text_after_skew.add_run()
		run2.text = text_after_skew_number
		skew_avg = text_frame.paragraphs[0]
		run3 = skew_avg.add_run()
		run3.text = avg_number
		font = run3.font
		font.bold = True
		text_agfter_avg = text_frame.paragraphs[0]
		run4 = text_agfter_avg.add_run()
		run4.text = text_after_avg_number
		segment = text_frame.paragraphs[0]
		run4 = segment.add_run()
		run4.text = seg

def round_string(stat):
	rounded = round(stat,1)
	string = str(rounded)
	return string

def va_text( 
		placeholder_index, 
		skew_number, 
		text_after_skew_number, 
		visit_alternative_competitor):
	text_frame = vo_slide.placeholders[placeholder_index].text_frame
	dem_type = vo_slide.placeholders[placeholder_index].text_frame
	skew = dem_type.paragraphs[0]
	run = skew.add_run()
	run.text = skew_number
	font = run.font
	font.bold = True
	# font.color.rgb = RGBColor(200,34,34)
	text_after_skew = text_frame.paragraphs[0]
	run2 = text_after_skew.add_run()
	run2.text = text_after_skew_number
	skew_avg = text_frame.paragraphs[0]
	run3 = skew_avg.add_run()
	run3.text = visit_alternative_competitor

#load necessary json data for competitors 
with open('kpi_payload_segment.json') as fp:
    kpis = json.load(fp)

Chains = ["Applebee's","Bahama Breeze Island Grille","Beef 'O' Brady's","BJ's Restaurant and Brewhouse","Bonefish Grill","Brio Tuscan Grille","Buffalo Wild Wings",
"California Pizza Kitchen","Carrabba's Italian Grill","Cheddar's Scratch Kitchen","The Cheesecake Factory","Chili's","Chuy's","Famous Dave's","Hooters","Joe’s Crab Shack",
"Logan's Roadhouse","LongHorn Steakhouse","Maggiano’s Little Italy","Mellow Mushroom","Miller's Ale House","Mimi's Cafe","O'Charley's","Olive Garden",
"On The Border Mexican Grill and Cantina","Outback Steakhouse","P.F. Chang's China Bistro","Pappasito's Cantina","Red Lobster","Red Robin Gourmet Burgers and Brews",
"Romano's Macaroni Grill","Ruby Tuesday","Seasons 52","Texas Roadhouse","TGI Fridays","Twin Peaks","Uno Pizzeria Grill","Yard House","Bob Evans", 
"Cracker Barrel Old Country Store", "Denny's", "First Watch","Friendly's","Huddle House","IHOP","Perkins Restaurants and Bakery","Shoney's","Village Inn","Waffle House",
"Au Bon Pain","Baja Fresh Mexican Grill","Blaze Pizza","Boston Market","Bruegger's Bagels","Chipotle Mexican Grill","Corner Bakery Cafe","Dickey's Barbecue Pit",
"Einstein Bros. Bagels","Fazoli's","Firehouse Subs","Five Guys Burgers and Fries","Fuddruckers","Habit Burger Grill, The","Jason's Deli","Jimmy John's Gourmet Sandwiches",
"McAlister's Deli", "Moe's Southwest Grill", "Newk's Eatery", "Noodles and Company", "Panda Express","Panera Bread","Pei Wei Asian Diner", "Pollo Campero","Pollo Tropical",
"Potbelly Sandwich Shop","Qdoba Mexican Eats", "Raising Cane's Chicken Fingers","Rubio's","Shake Shack","Smashburger","Which Wich Sandwiches", "Wingstop","Zaxby's","Zoes Kitchen",
"Arby's","Auntie Anne's","Baskin-Robbins","Ben and Jerry's","Bojangles' Famous Chicken 'N Biscuits","Burger King","Captain D's Seafood Kitchen","Caribou Coffee",
"Carl's Jr.","Carvel Ice Cream","Charley's Philly Steaks","Checkers Drive-In Restaurants","Chick-fil-A","Church's Chicken","Cicis","Cinnabon","Cold Stone Creamery","Culver's","Dairy Queen",
"Del Taco", "Domino's", "Dunkin' Donuts", "Godfather's Pizza","Golden Corral","Hardee's", "El Pollo Loco", "Freddy's Frozen Custard and Steakburger",
"HomeTown Buffet","In-N-Out Burger","Jack in the Box","Jamba Juice","Jet's Pizza","KFC","Krispy Kreme","Krystal Company","Little Caesars",
"Long John Silver's","Luby's","Marco's Pizza","McDonald's","Old Country Buffet","Papa John's","Papa Murphy's Pizza","Pinkberry","Pizza Hut","Popeyes Louisiana Kitchen",
"Quiznos","Sbarro","Sizzler","SONIC Drive-In","Starbucks","Steak 'n Shake","Subway","Taco Bell","Taco John's","Tim Hortons Cafe and Bake Shop","Tropical Smoothie Cafe",
"Wendy's","Whataburger","White Castle", "Hungry Howie's"]

for Chain in Chains:
	copyfile('input_conditional_2.pptx', DESTINATION_PATH +  Chain + '.pptx')
	Seg = kpis[slugify(Chain)]["Seg"]
	if Seg == "qsr":
		seg_dem_base = "72,630"
	elif Seg == "cdr":
		seg_dem_base = "36,498"
	elif Seg == "fsr":
		seg_dem_base = "42,390"
	else:
		seg_dem_base = "18,711"

	prs = Presentation('input_conditional_2.pptx')

	competitors = list(kpis[slugify(Chain)]["competitor_scores"].keys())
	competitor_series = competitors.append(slugify(Chain))
	competitor_series = competitors.append(Seg + "avg")

	competitor_values = list(kpis[slugify(Chain)]["competitor_scores"].values())
	competitor_tuple = list(zip(competitors[:-1], competitor_values))
	competitor_tuple.sort(key=lambda elem: elem[1])

	va_slide_competitor_list = [i[0] for i in competitor_tuple]
	va_slide_competitor_score_list = [i[1] for i in competitor_tuple]

	name_for_attributes = list()
	for c in competitors:
		name_for_attributes.append(kpis[c]["chain_name"])

	title_slide = prs.slides.add_slide(prs.slide_layouts[0])

	write_text(title_slide, 0, 'Top-Line Competitive Brand Assessment')
	write_text(title_slide, 11, 'KPI Stats')

	title_logo = title_slide.placeholders[19].insert_picture(LOGO_PATH + Chain + '.png') 
	title_logo.crop_top = 0
	title_logo.crop_bottom = 0
	title_logo.crop_right = 0
	title_logo.crop_left = 0
	mydate = datetime.datetime.now()
	current_date = str( mydate.strftime("%B"))	+ " " + str(mydate.strftime("%Y"))

	write_text(title_slide, 17, current_date)
	write_text(title_slide, 18, 'Powered by Consumer Brand Metrics')

	vo_slide = prs.slides.add_slide(prs.slide_layouts[1])
	if "'s" in Chain:
		vo_slide.placeholders[0].text = 'Top ' + Chain + " Competitors"
	elif Chain.endswith('s'):
		vo_slide.placeholders[0].text = 'Top ' + Chain + " Competitors"
	else:
		vo_slide.placeholders[0].text = 'Top ' + Chain + "'s Competitors"

	logo_pic(14,5)
	logo_pic(15,2)
	logo_pic(16,4)
	logo_pic(17,3)
	logo_pic(18,0)
	logo_pic(19,1)

	va1_text_params = [20, round_string((kpis[slugify(Chain)]["competitor_scores"][va_slide_competitor_list[5]]*100)) + "%", " of recent " + kpis[slugify(Chain)]["chain_name"] + " guests considered visiting ", kpis[va_slide_competitor_list[5]]["chain_name"]]
	va2_text_params = [21, round_string((kpis[slugify(Chain)]["competitor_scores"][va_slide_competitor_list[2]]*100)) + "%"," considered visiting ", kpis[va_slide_competitor_list[2]]["chain_name"]]
	va3_text_params = [22, round_string((kpis[slugify(Chain)]["competitor_scores"][va_slide_competitor_list[1]]*100)) + "%"," considered visiting ",kpis[va_slide_competitor_list[1]]["chain_name"]]
	va4_text_params = [23, round_string((kpis[slugify(Chain)]["competitor_scores"][va_slide_competitor_list[0]]*100)) + "%"," considered visiting ",kpis[va_slide_competitor_list[0]]["chain_name"]]
	va5_text_params = [24, round_string((kpis[slugify(Chain)]["competitor_scores"][va_slide_competitor_list[3]]*100)) + "%"," considered visiting ",kpis[va_slide_competitor_list[3]]["chain_name"]]
	va6_text_params = [25, round_string((kpis[slugify(Chain)]["competitor_scores"][va_slide_competitor_list[4]]*100)) + "%"," considered visiting ",kpis[va_slide_competitor_list[4]]["chain_name"]]
	va_text_lists = [va1_text_params,va2_text_params,va3_text_params,va4_text_params,va5_text_params,va6_text_params]

	for list_ in va_text_lists:
		print(len(list_))
		va_text(list_[0], list_[1], list_[2], list_[3])

	write_base_text(vo_slide, 
					26,
					"Total base: " 
					+ str(int(kpis[slugify(Chain)]["competitor_score_base"])) 
					+ " recent " 
					+ kpis[slugify(Chain)]["chain_name"] 
					+ " guests (Q3 - Q2'18)")
	
	write_rest_visit_text(vo_slide, 
							27, 
							round_string(kpis[slugify(Chain)]["Rest_visit"] * 100) 
							+ "%", 
							" would have gone to another restaurant as an alternative to " 
							+ Chain)
	
	gender = ['Male','Female']

	generation = ['Gen Z', 
				'Millennials', 
				'Generation X', 
				'Baby Boomers', 
				'Matures']

	ethnicity = ['Asian', 
				"Black/AA",	
				"Caucasian", 
				"Hispanic", 
				"Other"]

	household_income = [
						"< $25K",	
						"$25K - $35K",	
						"$35K - $50K", 
						"$50K - $75K", 
						"$75K - $100K", 
						"$100K - $150K", 
						"$150K +"
						]

	if kpis[slugify(Chain)]["attributes"]["Female"] > kpis[slugify(Chain)]["attributes"]["Male"]:
		gender = gender[1]
		gender_skew_for_Chain = kpis[slugify(Chain)]["attributes"]["Female"]
		gender_avg = kpis[Seg + "avg"]["attributes"]["Female"]
	else:
		gender = gender[0]
		gender_skew_for_Chain = kpis[slugify(Chain)]["attributes"]["Male"]
		gender_avg = kpis[Seg + "avg"]["attributes"]["Male"]

	generation_skew = max(kpis[slugify(Chain)]["attributes"]["Generation Z"], 
						kpis[slugify(Chain)]["attributes"]["Millennials"], 
						kpis[slugify(Chain)]["attributes"]["Generation X"], 
						kpis[slugify(Chain)]["attributes"]["Baby Boomers"], 
						kpis[slugify(Chain)]["attributes"]["Matures"])

	if generation_skew == kpis[slugify(Chain)]["attributes"]["Generation Z"]:
		generation = generation[0]
		generation_avg = kpis[Seg + "avg"]["attributes"]["Generation Z"]

	elif generation_skew == kpis[slugify(Chain)]["attributes"]["Millennials"]:
		generation = generation[1]
		generation_avg = kpis[Seg + "avg"]["attributes"]["Millennials"]

	elif generation_skew == kpis[slugify(Chain)]["attributes"]["Generation X"]:
		generation = generation[2]
		generation_avg = kpis[Seg + "avg"]["attributes"]["Generation X"]

	elif generation_skew == kpis[slugify(Chain)]["attributes"]["Baby Boomers"]:
		generation = generation[3]
		generation_avg = kpis[Seg + "avg"]["attributes"]["Baby Boomers"]

	else:
		generation = generation[4]
		generation_avg = kpis[Seg + "avg"]["attributes"]["Matures"]

	ethnic_skew = max(kpis[slugify(Chain)]["attributes"]["Asian"], 
					kpis[slugify(Chain)]["attributes"]["Black/AA"], 
					kpis[slugify(Chain)]["attributes"]["Caucasian"], 
					kpis[slugify(Chain)]["attributes"]["Hispanic"], 
					kpis[slugify(Chain)]["attributes"]["Other"])

	if ethnic_skew == kpis[slugify(Chain)]["attributes"]["Asian"]:
		ethnicity = ethnicity[0]
		ethnicity_avg = kpis[Seg + "avg"]["attributes"]["Asian"]

	elif ethnic_skew == kpis[slugify(Chain)]["attributes"]["Black/AA"]:
		ethnicity = ethnicity[1]
		ethnicity_avg = kpis[Seg + "avg"]["attributes"]["Black/AA"]

	elif ethnic_skew == kpis[slugify(Chain)]["attributes"]["Caucasian"]:
		ethnicity = ethnicity[2]
		ethnicity_avg = kpis[Seg + "avg"]["attributes"]["Caucasian"]

	elif ethnic_skew == kpis[slugify(Chain)]["attributes"]["Hispanic"]:
		ethnicity = ethnicity[3]
		ethnicity_avg = kpis[Seg + "avg"]["attributes"]["Hispanic"]

	else:
		ethnicity = ethnicity[4]
		ethnicity_avg = kpis[Seg + "avg"]["attributes"]["Other"]

	household_income_skew = max(kpis[slugify(Chain)]["attributes"]["Under $25,000"], 
								kpis[slugify(Chain)]["attributes"]["$25,000 - $34,999"], 
								kpis[slugify(Chain)]["attributes"]["$35,000 - $49,999"], 
								kpis[slugify(Chain)]["attributes"]["$50,000 - $74,999"], 
								kpis[slugify(Chain)]["attributes"]["$75,000 - $99,999"], 
								kpis[slugify(Chain)]["attributes"]["$100,000- $150,000"], 
								kpis[slugify(Chain)]["attributes"]["$150,000 +"])

	if household_income_skew == kpis[slugify(Chain)]["attributes"]["Under $25,000"]:
		household_income = household_income[0]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["Under $25,000"]

	elif household_income_skew == kpis[slugify(Chain)]["attributes"]["$25,000 - $34,999"]:
		household_income = household_income[1]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["$25,000 - $34,999"]

	elif household_income_skew == kpis[slugify(Chain)]["attributes"]["$35,000 - $49,999"]:
		household_income = household_income[2]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["$35,000 - $49,999"]

	elif household_income_skew == kpis[slugify(Chain)]["attributes"]["$50,000 - $74,999"]:
		household_income = household_income[3]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["$50,000 - $74,999"]

	elif household_income_skew == kpis[slugify(Chain)]["attributes"]["$75,000 - $99,999"]:
		household_income = household_income[4]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["$75,000 - $99,999"]

	elif household_income_skew == kpis[slugify(Chain)]["attributes"]["$100,000- $150,000"]:
		household_income = household_income[5]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["$100,000- $150,000"]

	else:
		household_income = household_income[6]
		hh_income_avg = kpis[Seg + "avg"]["attributes"]["$150,000 +"]

	dem_slide = prs.slides.add_slide(prs.slide_layouts[2])


	gender_text = ["gender", 18, round_string(gender_skew_for_Chain * 100) + "%", " Of " + Chain + "'s" + " frequent guest are " + gender + " compared to ", round_string(gender_avg * 100)+ "%", " across all ", Seg.upper() + "s"]
	generation_text = ["genderation", 19, round_string(generation_skew * 100) + "%", " Of frequent guest are " + generation + " compared to ", round_string(generation_avg * 100) + "%", " across all ", Seg.upper() + "s"]
	ethnicity_text = ["ethnicity", 20, round_string(ethnic_skew * 100) + "%", " Of frequent guest are " + ethnicity + " compared to ", round_string(ethnicity_avg * 100)+ "%", " across all ", Seg.upper() + "s"]
	hhi_text = ["household_income", 21, round_string(household_income_skew * 100) + "%", " Of frequent guest are " + household_income + " compared to ", round_string(hh_income_avg * 100)+ "%", " across all ", Seg.upper() + "s"]
	text_list = [gender_text, generation_text, ethnicity_text, hhi_text]

	for text in text_list:
			dem_text(text[0], text[1], text[2], text[3], text[4], text[5], text[6])
	
	write_text(dem_slide, 0, Chain + " Frequent Guest Demographic Skews")	
	write_text(dem_slide, 14, gender)
	write_text(dem_slide, 15, generation)
	write_text(dem_slide, 16, ethnicity)
	write_text(dem_slide, 17, household_income)
	write_base_text(dem_slide, 
					22, 
					"Frequent guest = consumers that visit the chain once a month or more Base: " 
					+ seg_dem_base 
					+ " once a month+ fast casual consumers (Q3 – Q2‘18)")
	
	if Seg == 'msr':
		most_important_attribute_slide = prs.slides.add_slide(prs.slide_layouts[8])

	elif Seg == 'cdr':
		most_important_attribute_slide = prs.slides.add_slide(prs.slide_layouts[9])

	elif Seg == 'fcr':
		most_important_attribute_slide = prs.slides.add_slide(prs.slide_layouts[10])

	else:
		most_important_attribute_slide = prs.slides.add_slide(prs.slide_layouts[11])

	#FOOD ATTRIBUTES SLIDE
	if Seg == 'qsr':
		food_slide = prs.slides.add_slide(prs.slide_layouts[16])
		write_text(food_slide, 0, 'Food Attributes')
		write_text(food_slide, 17, 'Food Taste & Flavor')
		write_text(food_slide, 18, 'Food Quality')
		write_text(food_slide, 19, 'Food Quality When Ordered for Takeout*')
		write_text(food_slide, 
					20, 
					"Based on your recent visit, how would you rate the chain on _________?" 
					+ "                                      "
					+ "     Total base: 700 recent guests per brand (Q3 – Q2‘18) Showing percentage selecting “very good” (top-box rating)")

	else:
		food_slide = prs.slides.add_slide(prs.slide_layouts[13])
		write_text(food_slide, 0, 'Food Attributes')	
		write_text(food_slide, 18, 'Food Taste & Flavor')
		write_text(food_slide, 19, 'Food Quality')
		write_text(food_slide, 
					20, 
					"Based on your recent visit, how would you rate the chain on _________?" 
					+ "                                      "
					+ "     Total base: 700 recent guests per brand (Q3 – Q2‘18) Showing percentage selecting “very good” (top-box rating)")

	food_quality_series = (kpis[competitors[0]]["attributes"]["Food quality"], 
							kpis[competitors[1]]["attributes"]["Food quality"], 
							kpis[competitors[2]]["attributes"]["Food quality"], 
							kpis[competitors[3]]["attributes"]["Food quality"], 
							kpis[competitors[4]]["attributes"]["Food quality"], 
							kpis[competitors[5]]["attributes"]["Food quality"],
							kpis[slugify(Chain)]["attributes"]["Food quality"], 
							kpis[Seg + "avg"]["attributes"]["Food quality"])

	taste_flavor_series = (kpis[competitors[0]]["attributes"]["Food taste and flavor"], 
							kpis[competitors[1]]["attributes"]["Food taste and flavor"], 
							kpis[competitors[2]]["attributes"]["Food taste and flavor"], 
							kpis[competitors[3]]["attributes"]["Food taste and flavor"], 
							kpis[competitors[4]]["attributes"]["Food taste and flavor"], 
							kpis[competitors[5]]["attributes"]["Food taste and flavor"],
							kpis[slugify(Chain)]["attributes"]["Food taste and flavor"], 
							kpis[Seg + "avg"]["attributes"]["Food taste and flavor"])

	takeout_food_quality_series =(kpis[competitors[0]]["attributes"]["Food quality takeout"],
							kpis[competitors[1]]["attributes"]["Food quality takeout"],
							kpis[competitors[2]]["attributes"]["Food quality takeout"],
							kpis[competitors[3]]["attributes"]["Food quality takeout"],
							kpis[competitors[4]]["attributes"]["Food quality takeout"],
							kpis[competitors[5]]["attributes"]["Food quality takeout"],
							kpis[slugify(Chain)]["attributes"]["Food quality takeout"],
							kpis[Seg + "avg"]["attributes"]["Food quality takeout"])


	create_chart('food_quality', 
				food_slide, 
				15, 
				name_for_attributes, 
				food_quality_series, 
				"Food quality")

	create_chart('taste_flavor', 
				food_slide, 
				14, 
				name_for_attributes, 
				taste_flavor_series, 
				"Food taste and flavor")

	if Seg == 'qsr':
		create_chart('taste_flavor', 
					food_slide, 
					16, 
					name_for_attributes, 
					takeout_food_quality_series, 
					"Food quality takeout")

	crave_slide = prs.slides.add_slide(prs.slide_layouts[4])
	craveable_keys = list(kpis[slugify(Chain)]["craveable_items"].keys())
	craveable_series = list(kpis[slugify(Chain)]["craveable_items"].values())

	crave_chart("craveable_items", 
				crave_slide, 
				14, 
				craveable_keys[:4], 
				craveable_series[:4])

	write_text(crave_slide, 0, 'Most Craveable ' + Chain + " Items")

	write_text(crave_slide, 
					15, 
				str(kpis[slugify(Chain)]["craveable_items_base"]).strip('.0') 
				+ ' recent ' 
				+ Chain 
				+ ' guests who rated the chain as "good" or "very good" for craveable items"')

	if Seg == 'qsr':
		acc_slide = prs.slides.add_slide(prs.slide_layouts[19])
	elif Seg == "fcr":
		acc_slide = prs.slides.add_slide(prs.slide_layouts[18])
	else:
		acc_slide = prs.slides.add_slide(prs.slide_layouts[20])

	write_text(acc_slide, 0, 'Accuracy \nAttribute')
	write_text(acc_slide, 19, 'Order Accuracy')

	write_text(acc_slide, 
				18, 
				"Based on your recent visit, how would you rate the chain on _________?" 
				+ "                                      "
				+ "     Total base: 700 recent guests per brand (Q3 – Q2‘18) Showing percentage selecting “very good” (top-box rating)")	
	# write_text(acc_slide, 0, 'Accuracy Attribute')

	order_accuracy_series = (kpis[competitors[0]]["attributes"]["Order accuracy"], 
							kpis[competitors[1]]["attributes"]["Order accuracy"], 
							kpis[competitors[2]]["attributes"]["Order accuracy"], 
							kpis[competitors[3]]["attributes"]["Order accuracy"], 
							kpis[competitors[4]]["attributes"]["Order accuracy"], 
							kpis[competitors[5]]["attributes"]["Order accuracy"],
							kpis[slugify(Chain)]["attributes"]["Order accuracy"], 
							kpis[Seg + "avg"]["attributes"]["Order accuracy"])

	create_chart("order_accuracy", 
					acc_slide, 
					17, 
					name_for_attributes, 
					order_accuracy_series, 
					"Order accuracy")
	
	int_series = (kpis[competitors[0]]["attributes"]["Interior cleanliness"], 
				kpis[competitors[1]]["attributes"]["Interior cleanliness"], 
				kpis[competitors[2]]["attributes"]["Interior cleanliness"], 
				kpis[competitors[3]]["attributes"]["Interior cleanliness"], 
				kpis[competitors[4]]["attributes"]["Interior cleanliness"], 
				kpis[competitors[5]]["attributes"]["Interior cleanliness"],
				kpis[slugify(Chain)]["attributes"]["Interior cleanliness"], 
				kpis[Seg + "avg"]["attributes"]["Interior cleanliness"])

	kitchenpreparea_series = (kpis[competitors[0]]["attributes"]["Kitchen/food prep area cleanliness"], 
							kpis[competitors[1]]["attributes"]["Kitchen/food prep area cleanliness"], 
							kpis[competitors[2]]["attributes"]["Kitchen/food prep area cleanliness"], 
							kpis[competitors[3]]["attributes"]["Kitchen/food prep area cleanliness"], 
							kpis[competitors[4]]["attributes"]["Kitchen/food prep area cleanliness"], 
							kpis[competitors[5]]["attributes"]["Kitchen/food prep area cleanliness"],
							kpis[slugify(Chain)]["attributes"]["Kitchen/food prep area cleanliness"], 
							kpis[Seg + "avg"]["attributes"]["Kitchen/food prep area cleanliness"])

	dish_series = (kpis[competitors[0]]["attributes"]["Dishware/glassware/silverware cleanliness"], 
				kpis[competitors[1]]["attributes"]["Dishware/glassware/silverware cleanliness"], 
				kpis[competitors[2]]["attributes"]["Dishware/glassware/silverware cleanliness"], 
				kpis[competitors[3]]["attributes"]["Dishware/glassware/silverware cleanliness"], 
				kpis[competitors[4]]["attributes"]["Dishware/glassware/silverware cleanliness"], 
				kpis[competitors[5]]["attributes"]["Dishware/glassware/silverware cleanliness"],
				kpis[slugify(Chain)]["attributes"]["Dishware/glassware/silverware cleanliness"], 
				kpis[Seg + "avg"]["attributes"]["Dishware/glassware/silverware cleanliness"])

	if Seg == 'qsr':
		clean_slide = prs.slides.add_slide(prs.slide_layouts[17])

		write_text(clean_slide, 0, "Cleanliness Attributes")

		create_chart('kitchenpreparea_chart', 
					clean_slide, 
					15, 
					name_for_attributes, 
					kitchenpreparea_series, 
					"Kitchen/food prep area cleanliness")

		create_chart('int_chart', 
					clean_slide, 
					14, 
					name_for_attributes, 
					int_series, 
					"Interior cleanliness")

		write_text(clean_slide, 18, 'Interior Cleanliness')
		write_text(clean_slide, 19, 'Kitchen or Food Prep Area Cleanliness')

	elif Seg == "fcr":

		clean_slide = prs.slides.add_slide(prs.slide_layouts[14])

		write_text(clean_slide, 0, "Cleanliness Attributes")

		dish_params = ['dish_chart', clean_slide, 16, name_for_attributes, dish_series, "Dishware/glassware/silverware cleanliness"]
		interior_params = ['int_chart', clean_slide, 14, name_for_attributes, int_series, "Interior cleanliness"]
		kitchen_params = ['kitchenpreparea_chart', clean_slide, 15, name_for_attributes, kitchenpreparea_series, "Kitchen/food prep area cleanliness"]
		param_list = [dish_params, interior_params, kitchen_params]

		for list_ in param_list:
			create_chart(list_[0],list_[1],list_[2],list_[3],list_[4],list_[5])

		write_text(clean_slide, 17, 'Cleanliness of Dishware and Silverware')
		write_text(clean_slide, 18, 'Interior Cleanliness')
		write_text(clean_slide, 19, 'Kitchen or Food Prep Area Cleanliness')

	else:
		clean_slide = prs.slides.add_slide(prs.slide_layouts[15])

		write_text(clean_slide, 0, "Cleanliness Attributes")

		dish_params = ['dish_chart', clean_slide, 16, name_for_attributes, dish_series, "Dishware/glassware/silverware cleanliness"]
		interior_params = ['int_chart', clean_slide, 14, name_for_attributes, int_series, "Interior cleanliness"]
		kitchen_params = ['kitchenpreparea_chart', clean_slide, 15, name_for_attributes, kitchenpreparea_series, "Kitchen/food prep area cleanliness"]
		param_list = [dish_params, interior_params, kitchen_params]

		for list_ in param_list:
			create_chart(list_[0],list_[1],list_[2],list_[3],list_[4],list_[5])

		write_text(clean_slide, 17, 'Cleanliness of Dishware and Silverware')
		write_text(clean_slide, 18, 'Interior Cleanliness')
		write_text(clean_slide, 19, 'Kitchen or Food Prep Area Cleanliness')
	
	write_text(clean_slide, 
				20, 
				"Based on your recent visit, how would you rate the chain on _________?" 
				+ "                                      "
				+ "     Total base: 700 recent guests per brand (Q3 – Q2‘18) Showing percentage selecting “very good” (top-box rating)")
	
	overall_slide = prs.slides.add_slide(prs.slide_layouts[7])
	
	write_text(overall_slide, 
				13, 
				"Based on your recent visit, how would you rate the chain on _________?" 
				+ "                                      "
				+ "     Total base: 700 recent guests per brand (Q3 – Q2‘18) Showing percentage selecting “very good” (top-box rating)")

	write_text(overall_slide, 0, 'Overall Visit Satisfaction')

	overall_series = (kpis[competitors[0]]["attributes"]["Overall Rating"], 
					kpis[competitors[1]]["attributes"]["Overall Rating"], 
					kpis[competitors[2]]["attributes"]["Overall Rating"], 
					kpis[competitors[3]]["attributes"]["Overall Rating"], 
					kpis[competitors[4]]["attributes"]["Overall Rating"], 
					kpis[competitors[5]]["attributes"]["Overall Rating"],
					kpis[slugify(Chain)]["attributes"]["Overall Rating"], 
					kpis[Seg + "avg"]["attributes"]["Overall Rating"])

	create_chart('overall_chart', 
					overall_slide, 
					12, 
					name_for_attributes, 
					overall_series, 
					"Overall Rating")

	prs.save(DESTINATION_PATH + Chain + '.pptx')